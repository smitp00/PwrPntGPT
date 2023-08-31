from django.http import HttpResponse
from django.shortcuts import render, redirect
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from .gpt_service import get_summary_from_gpt  # Place gpt_service.py in the my_app directory

# Combined upload view for handling both PowerPoint and Word uploads
def upload(request):
    if request.method == 'POST':
        ppt_file = request.FILES.get('ppt_file')
        word_file = request.FILES.get('word_file')

        # Handle PowerPoint upload
        if ppt_file and ppt_file.name.endswith('.pptx'):
            try:
                prs = Presentation(ppt_file)
                slides_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slides_text.append(run.text)
                        # Handle table shape
                        if shape.shape_type == 19:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    slides_text.append(cell.text)
                if not slides_text:
                    return render(request, 'error.html', {'error_message': 'The presentation appears to be empty.'})
                summarized_content = get_summary_from_gpt(' '.join(slides_text))
                return render(request, 'summaryppt.html', {'original_content': slides_text, 'summarized_content': summarized_content})
            except Exception as e:
                return render(request, 'error.html', {'error_message': str(e)})

        # Handle Word upload
        elif word_file and word_file.name.endswith('.docx'):
            document = Document(word_file)
            doc_text = []
            for para in document.paragraphs:
                doc_text.append(para.text)
            # Combine all the text into a single string
            full_text = '\n'.join(doc_text)
            # Get the summarized content
            summary = get_summary_from_gpt(full_text)
            return render(request, 'summaryword.html', {'summary': summary})
            
    return render(request, 'upload.html')